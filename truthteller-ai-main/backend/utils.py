"""
utils.py
────────
File parsing, feature extraction, and text analysis helpers.

FIXES vs original:
  • Replaced PyPDF2 (abandoned) with pypdf (actively maintained)
  • Removed calls to load_ml_model() / load_simple_model() which
    were undefined in this file (they live in app.py)
  • predict_comprehensive() now uses the stored LR weights from the
    JSON model instead of hard-coded rule thresholds
  • analyze_text() cascades cleanly: real_model → fallback
  • The model is now a true Scikit-Learn pipeline (TF-IDF + LR)
"""

import os
import json
import math
import re
import numpy as np
from collections import Counter

# pypdf replaces the abandoned PyPDF2 — drop-in compatible API
try:
    from pypdf import PdfReader
except ImportError:
    try:
        import PyPDF2
        class PdfReader:   # shim for legacy installs
            def __init__(self, f):
                self._r = PyPDF2.PdfReader(f)
            @property
            def pages(self):
                return self._r.pages
    except ImportError:
        PdfReader = None

from docx import Document
from pptx import Presentation
from werkzeug.utils import secure_filename
import joblib

import nltk
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.corpus import stopwords

# ─── NLTK data ────────────────────────────────────────────────────────────────
for _resource, _path in [
    ("tokenizers/punkt", "punkt"),
    ("corpora/stopwords", "stopwords"),
    ("tokenizers/punkt_tab", "punkt_tab"),
]:
    try:
        nltk.data.find(_resource)
    except LookupError:
        try:
            nltk.download(_path, quiet=True)
        except Exception:
            pass

# ─── Constants ────────────────────────────────────────────────────────────────
ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx"}



# ─── Model loading ────────────────────────────────────────────────────────────
REAL_MODEL_PATH = os.path.join(
    os.path.dirname(__file__), "models", "real_model.pkl"
)

real_model = None


def load_real_model():
    """Load the real ML model pipeline from PKL."""
    global real_model
    try:
        if os.path.exists(REAL_MODEL_PATH):
            real_model = joblib.load(REAL_MODEL_PATH)
            print("[OK] Real model loaded successfully")
            return True
        else:
            print("[WARN] Real model not found - run train_real_model.py first")
            return False
    except Exception as e:
        print(f"[ERROR] Error loading real model: {e}")
        return False


# Load on import
load_real_model()


# ─── File helpers ─────────────────────────────────────────────────────────────
def allowed_file(filename: str) -> bool:
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    if PdfReader is None:
        print("No PDF library available. Install: pip install pypdf")
        return text
    try:
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
    except Exception as e:
        print(f"Error extracting PDF: {e}")
    return text


def extract_text_from_docx(file_path: str) -> str:
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        print(f"Error extracting DOCX: {e}")
    return text


def extract_text_from_pptx(file_path: str) -> str:
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting PPTX: {e}")
    return text


def extract_text_from_file(file_path: str, file_extension: str) -> str:
    ext = file_extension.lower().lstrip(".")
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext == "docx":
        return extract_text_from_docx(file_path)
    elif ext == "pptx":
        return extract_text_from_pptx(file_path)
    return ""







# ─── Main public API ──────────────────────────────────────────────────────────
def analyze_text(text: str) -> dict:
    """
    Analyse text and return {"prediction": "AI"|"Human", "confidence": float}.

    Priority:
      1. real_model (TF-IDF + LR pipeline)
      2. fallback heuristics
    """
    if not text or not text.strip():
        return {"prediction": "Unknown", "confidence": 0.0}

    if real_model is not None:
        try:
            # Predict
            prob_ai = float(real_model.predict_proba([text])[0][1])
            pred = "AI" if prob_ai >= 0.5 else "Human"
            conf = max(prob_ai, 1.0 - prob_ai)
            return {"prediction": pred, "confidence": round(conf, 3)}

        except Exception as e:
            print(f"Real model inference error: {e}")

    # Fallback
    return _analyze_text_fallback(text)





def _analyze_text_fallback(text: str) -> dict:
    """Heuristic fallback when no model is available."""
    words = re.findall(r"\b\w+\b", text.lower())
    n     = max(len(words), 1)

    # AI indicators
    word_freq    = Counter(words)
    lex_div      = len(word_freq) / n
    bigrams      = [f"{words[i]} {words[i+1]}" for i in range(n - 1)]
    bg_rep       = 1 - (len(set(bigrams)) / len(bigrams)) if bigrams else 0.0
    entropy      = -sum((c / n) * math.log2(c / n) for c in word_freq.values() if c > 0)
    perplexity   = 2 ** entropy

    ai_score = 0.0
    if lex_div < 0.4:    ai_score += 0.2
    if bg_rep  > 0.1:    ai_score += 0.2
    if perplexity < 50:  ai_score += 0.2

    prob_ai  = 0.5 + ai_score * 0.5
    pred     = "AI" if prob_ai >= 0.5 else "Human"
    conf     = round(min(max(prob_ai, 1.0 - prob_ai), 0.95), 3)
    return {"prediction": pred, "confidence": conf}